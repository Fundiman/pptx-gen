import subprocess
import os
import argparse
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import requests

def fetch_images(topic, count):
    # Run gimagesearch to fetch images
    command = f'python3 gimagesearch.py "{topic}" {count}'
    subprocess.run(command, shell=True, check=True)

    # Collect the image filenames
    image_files = [f for f in os.listdir() if f.startswith(f"{topic}_") and f.endswith(('.jpeg', '.png'))]
    return image_files

def fetch_topic_data(topic):
    search_url = f'https://en.wikipedia.org/wiki/{topic.replace(" ", "_")}'
    response = requests.get(search_url)
    soup = BeautifulSoup(response.text, 'html.parser')
    
    paragraphs = soup.find_all('p')
    summary = []
    
    for paragraph in paragraphs:
        text = paragraph.get_text().strip()
        if text and len(text) > 100:
            summary.append(text)
        if len(summary) >= 5:
            break
    
    return summary

def random_color():
    return RGBColor(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))

def random_font():
    fonts = ['Arial', 'Calibri', 'Times New Roman', 'Verdana', 'Tahoma']
    return random.choice(fonts)

def add_title_slide(prs, topic, bg_color):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = topic
    subtitle.text = f'Exploring and Learning about {topic}'
    
    title_paragraph = title.text_frame.paragraphs[0]
    subtitle_paragraph = subtitle.text_frame.paragraphs[0]
    
    title_paragraph.font.size = Pt(48)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = random_color()
    title_paragraph.font.name = random_font()
    
    subtitle_paragraph.font.size = Pt(24)
    subtitle_paragraph.font.color.rgb = random_color()
    subtitle_paragraph.font.name = random_font()
    
    # Set background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

def add_combined_slide(prs, title_text, content_text, image_path, bg_color):
    slide_layout = prs.slide_layouts[5]  # Title and Content Layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = random_color()
    title_paragraph.font.name = random_font()
    
    # Set content text
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = content_text[:1000]  # Limit text to avoid overflow
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = random_color()
        paragraph.font.name = random_font()
    content_frame.word_wrap = True
    
    # Set image
    left = Inches(6)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(4))
    
    # Set background color
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    slide.shapes._spTree.insert(0, shape._element)

def delete_all_images():
    for file in os.listdir():
        if file.lower().endswith(('.jpeg', '.png')):
            try:
                os.remove(file)
                print(f"Deleted {file}")
            except Exception as e:
                print(f"Error deleting {file}: {e}")

def create_presentation(topic, image_files, paragraphs, slide_count, bg_color):
    prs = Presentation()
    
    add_title_slide(prs, topic, bg_color)
    
    text_titles = [
        f"What is {topic}?",
        f"History of {topic}",
        f"Key Facts About {topic}",
        f"Significance of {topic}",
        f"Interesting Aspects of {topic}"
    ]
    
    slide_index = 0
    for i, para in enumerate(paragraphs):
        if slide_index >= slide_count - 1:
            break
        
        # Use images cyclically for each text slide
        image_file = image_files[i % len(image_files)]
        slide_title = text_titles[slide_index % len(text_titles)]
        
        add_combined_slide(prs, slide_title, para, image_file, bg_color)
        slide_index += 1
    
    file_name = f'{topic}_presentation.pptx'
    file_path = os.path.join(os.getcwd(), file_name)
    prs.save(file_path)
    print(f'Presentation saved as {file_path}')
    
    # Delete all image files after creating the presentation
    delete_all_images()

def main():
    parser = argparse.ArgumentParser(description='Generate a PowerPoint presentation with images and text.')
    parser.add_argument('--topic', type=str, required=True, help='Topic for the presentation')
    parser.add_argument('--slides', type=int, default=10, help='Number of slides for the presentation')
    parser.add_argument('--images', type=int, default=10, help='Number of images to fetch')
    parser.add_argument('--bg_color', type=str, default='white', choices=['white', 'black'], help='Background color of the slides (white or black)')
    
    args = parser.parse_args()
    
    topic = args.topic
    slide_count = args.slides
    image_count = args.images
    bg_color = RGBColor(255, 255, 255) if args.bg_color == 'white' else RGBColor(0, 0, 0)
    
    print("Fetching images...")
    fetch_images(topic, image_count)
    
    image_files = [f for f in os.listdir() if f.startswith(f"{topic}_") and f.endswith(('.jpeg', '.png'))]
    
    if not image_files:
        print("No images found. Ensure gimagesearch.py is working correctly.")
        return
    
    print("Fetching data...")
    paragraphs = fetch_topic_data(topic)
    
    if len(paragraphs) < 2:
        print("Not enough data retrieved. Try a different topic.")
        return
    
    print("Creating presentation...")
    create_presentation(topic, image_files, paragraphs, slide_count, bg_color)
    
    print("Presentation created successfully!")

if __name__ == "__main__":
    main()

